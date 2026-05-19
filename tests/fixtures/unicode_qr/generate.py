#!/usr/bin/env python3
"""Generate a suite of test files each carrying a Unicode-block-art QR code
in their text body. Each file encodes a distinct URL so we can confirm in
the results which decode came from which format."""
import os
import qrcode
import qrcode.constants
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

OUT = "/tmp/qr_fixtures"
os.makedirs(OUT, exist_ok=True)

def render_qr_art(payload: str) -> str:
    qr = qrcode.QRCode(version=2,
                       error_correction=qrcode.constants.ERROR_CORRECT_M,
                       box_size=1, border=2)
    qr.add_data(payload)
    qr.make(fit=True)
    m = qr.get_matrix()
    lines = []
    for y in range(0, len(m), 2):
        row = []
        for x in range(len(m[0])):
            t = m[y][x]
            b = m[y + 1][x] if y + 1 < len(m) else False
            if t and b:       row.append('█')
            elif t and not b: row.append('▀')
            elif not t and b: row.append('▄')
            else:             row.append(' ')
        lines.append(''.join(row))
    return '\n'.join(lines)

def fold_ics(text: str) -> str:
    # ICS single-line escape: \n for inline newline, line-fold any line >75 chars
    s = text.replace('\n', '\\n')
    out = []
    while len(s) > 74:
        out.append(s[:74])
        s = ' ' + s[74:]
    out.append(s)
    return '\r\n'.join(out)

# ---- ICS ----
ics_url = "https://test.example/qr-ics"
ics_art = render_qr_art(ics_url)
ics = f"""BEGIN:VCALENDAR\r
VERSION:2.0\r
PRODID:-//qr fixture//EN\r
METHOD:REQUEST\r
BEGIN:VEVENT\r
UID:qr-fixture-ics\r
DTSTAMP:20251101T120000Z\r
DTSTART:20251115T140000Z\r
DTEND:20251115T150000Z\r
SUMMARY:Scan this\r
DESCRIPTION:{fold_ics('Scan the QR below:\n' + ics_art)}\r
END:VEVENT\r
END:VCALENDAR\r
"""
with open(f"{OUT}/qr_fixture.ics", "w", encoding="utf-8") as f: f.write(ics)

# ---- TXT ----
txt_url = "https://test.example/qr-txt"
with open(f"{OUT}/qr_fixture.txt", "w", encoding="utf-8") as f:
    f.write(f"Scan the QR below to win:\n\n{render_qr_art(txt_url)}\n")

# ---- HTML (QR inside <pre>) ----
html_url = "https://test.example/qr-html"
html_art = render_qr_art(html_url)
with open(f"{OUT}/qr_fixture.html", "w", encoding="utf-8") as f:
    f.write(f"""<!doctype html>
<html><head><meta charset="utf-8"><title>QR test</title></head>
<body>
<p>Scan the code:</p>
<pre style="font-family: monospace; line-height: 1;">{html_art}</pre>
</body></html>
""")

# ---- EML (multipart with plaintext QR + HTML QR) ----
eml_url = "https://test.example/qr-eml"
eml_art = render_qr_art(eml_url)
msg = MIMEMultipart("alternative")
msg["Subject"] = "Important — please scan"
msg["From"] = "phisher@evil.example"
msg["To"] = "victim@target.example"
plain = MIMEText("Scan the QR below:\n\n" + eml_art, "plain", "utf-8")
html_body = MIMEText(
    f"<html><body><p>Scan:</p><pre>{eml_art}</pre></body></html>", "html", "utf-8")
msg.attach(plain)
msg.attach(html_body)
with open(f"{OUT}/qr_fixture.eml", "wb") as f:
    f.write(msg.as_bytes())

# ---- RTF (QR in plain text) ----
rtf_url = "https://test.example/qr-rtf"
rtf_art = render_qr_art(rtf_url)
# Escape unicode for RTF — chars > 0x7F become \\uNNNNN?
def rtf_escape(s):
    out = []
    for ch in s:
        c = ord(ch)
        if ch == '\\' or ch == '{' or ch == '}':
            out.append('\\' + ch)
        elif ch == '\n':
            out.append('\\par\n')
        elif c < 0x80:
            out.append(ch)
        else:
            # signed 16-bit value
            v = c if c < 0x8000 else c - 0x10000
            out.append(f"\\u{v}?")
    return ''.join(out)

with open(f"{OUT}/qr_fixture.rtf", "w", encoding="utf-8") as f:
    f.write("{\\rtf1\\ansi\\ansicpg1252\\uc1\n")
    f.write("{\\fonttbl\\f0\\fcourier Courier;}\n")
    f.write("\\f0\\fs16\n")
    f.write("Scan the QR below:\\par\n")
    f.write(rtf_escape(rtf_art))
    f.write("\n}")

# ---- DOCX (Word body with QR in monospace) ----
docx_url = "https://test.example/qr-docx"
try:
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    p = doc.add_paragraph()
    r = p.add_run("Scan the QR below to confirm your account:")
    doc.add_paragraph()
    p2 = doc.add_paragraph()
    r2 = p2.add_run(render_qr_art(docx_url))
    r2.font.name = 'Courier New'
    r2.font.size = Pt(6)
    doc.save(f"{OUT}/qr_fixture.docx")
    print("docx ok")
except Exception as e:
    print(f"docx error: {e}")

# ---- PPTX (slide notes / text frame) ----
pptx_url = "https://test.example/qr-pptx"
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.placeholders[0].text_frame
    tf.text = "Scan QR to confirm"
    # Add a text box with QR art in monospace
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(6))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = render_qr_art(pptx_url)
    r.font.name = 'Courier New'
    r.font.size = Pt(6)
    prs.save(f"{OUT}/qr_fixture.pptx")
    print("pptx ok")
except Exception as e:
    print(f"pptx error: {e}")

# ---- PDF (QR art rendered with FPDF) ----
pdf_url = "https://test.example/qr-pdf"
try:
    from fpdf import FPDF
    pdf = FPDF(orientation='P', unit='mm', format='A4')
    pdf.add_page()
    pdf.set_font("Courier", size=8)
    pdf.cell(0, 10, "Scan the QR below:", ln=1)
    pdf.set_font("Courier", size=4)
    for line in render_qr_art(pdf_url).split('\n'):
        pdf.cell(0, 1.5, line, ln=1)
    pdf.output(f"{OUT}/qr_fixture.pdf")
    print("pdf ok")
except Exception as e:
    print(f"pdf error: {e}")

print()
print(f"Fixtures in {OUT}:")
for f in sorted(os.listdir(OUT)):
    print(f"  {f}  ({os.path.getsize(os.path.join(OUT, f))} bytes)")
